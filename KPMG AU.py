#!/usr/bin/env python
# coding: utf-8

# In[1]:


import pandas as pd

# Load the Excel file
file_path = 'KPMG_VI_New_raw_data_update_final.xlsx'
excel_data = pd.ExcelFile(file_path)

# Display sheet names to understand the structure of the file
sheet_names = excel_data.sheet_names
sheet_names


# In[2]:


# Load and display the first few rows of each sheet

# Load the data from each sheet
title_sheet_df = pd.read_excel(file_path, sheet_name='Title Sheet')
transactions_df = pd.read_excel(file_path, sheet_name='Transactions')
new_customer_list_df = pd.read_excel(file_path, sheet_name='NewCustomerList')
customer_demographic_df = pd.read_excel(file_path, sheet_name='CustomerDemographic')
customer_address_df = pd.read_excel(file_path, sheet_name='CustomerAddress')

# Display the first few rows of each sheet to get an overview
title_sheet_df.head(), transactions_df.head(), new_customer_list_df.head(), customer_demographic_df.head(), customer_address_df.head()


# In[4]:


# Transactions Sheet Analysis: First, we'll check for missing values, duplicates, and data format consistency in the Transactions sheet. 

# Check for missing values, duplicates, and data format consistency in the Transactions sheet

# Check for missing values
missing_values_transactions = transactions_df.isnull().sum()

# Check for duplicates
duplicate_transactions = transactions_df.duplicated().sum()

# Display the results
missing_values_transactions, duplicate_transactions


# In[ ]:


# The Transactions sheet has missing values in multiple columns. Specifically, columns with missing values include Unnamed: 4 and Unnamed: 6 to Unnamed: 12. There are no duplicates in the Transactions sheet.


# In[5]:


# New Customer List Analysis: Next, let's check for missing values, duplicates, and data format consistency in the NewCustomerList sheet.

# Check for missing values, duplicates, and data format consistency in the New Customer List sheet

# Check for missing values
missing_values_new_customer_list = new_customer_list_df.isnull().sum()

# Check for duplicates
duplicate_new_customer_list = new_customer_list_df.duplicated().sum()

# Display the results
missing_values_new_customer_list, duplicate_new_customer_list


# In[ ]:


# The NewCustomerList sheet has missing values in several columns. Specifically, columns with significant missing values include Unnamed: 1, Unnamed: 4, Unnamed: 5, and Unnamed: 6. There are no duplicates in the NewCustomerList sheet.


# In[7]:


# Customer Demographic Analysis: Now, let's check for missing values, duplicates, and data format consistency in the CustomerDemographic sheet.

# Check for missing values, duplicates, and data format consistency in the Customer Demographic sheet

# Check for missing values
missing_values_customer_demographic = customer_demographic_df.isnull().sum()

# Check for duplicates
duplicate_customer_demographic = customer_demographic_df.duplicated().sum()

# Display the results
missing_values_customer_demographic, duplicate_customer_demographic


# In[ ]:


# The CustomerDemographic sheet has missing values in several columns, notably Unnamed: 2, Unnamed: 5, Unnamed: 6, Unnamed: 7, Unnamed: 10, and Unnamed: 12. There are no duplicates in the CustomerDemographic sheet.


# In[8]:


# Customer Address Analysis: Finally, let's check for missing values, duplicates, and data format consistency in the CustomerAddress sheet.

# Check for missing values, duplicates, and data format consistency in the Customer Address sheet

# Check for missing values
missing_values_customer_address = customer_address_df.isnull().sum()

# Check for duplicates
duplicate_customer_address = customer_address_df.duplicated().sum()

# Display the results
missing_values_customer_address, duplicate_customer_address


# In[ ]:


# The CustomerAddress sheet has no missing values and no duplicates.


# In[12]:


get_ipython().system('pip install python-pptx')


# In[13]:


# Preparing the Presentation

from pptx import Presentation
from pptx.util import Inches

# Load the template presentation
ppt_template_path = 'Module_2_Template_slide.pptx'
presentation = Presentation(ppt_template_path)

# Define content for each slide
slides_content = {
    'Introduction': {
        'title': "Introduction",
        'content': [
            "Sprocket Central Pty Ltd - Data Analytics Approach",
            "Team: [Division Name], [Engagement Manager], [Senior Consultant], [Junior Consultant]",
            "Purpose: Outline the approach for analyzing customer data to drive business value."
        ]
    },
    'Agenda': {
        'title': "Agenda",
        'content': [
            "Introduction",
            "Data Exploration",
            "Model Development",
            "Interpretation"
        ]
    },
    'Data Exploration': {
        'title': "Data Exploration",
        'content': [
            "Objective: Understand the data distributions and identify data quality issues.",
            "Activities:",
            "- Data Loading and Initial Inspection",
            "- Missing Value Analysis and Imputation Strategies",
            "- Data Consistency Checks and Cleaning",
            "- Exploratory Data Analysis (EDA)",
            "- Feature Engineering (e.g., converting D.O.B to age or age groups)"
        ]
    },
    'Model Development': {
        'title': "Model Development",
        'content': [
            "Objective: Develop predictive models to identify high-value customers.",
            "Activities:",
            "- Data Transformation and Scaling",
            "- Feature Selection and Importance Analysis",
            "- Model Selection (e.g., Logistic Regression, Decision Trees, Random Forest)",
            "- Model Training and Validation",
            "- Hyperparameter Tuning"
        ]
    },
    'Interpretation': {
        'title': "Interpretation",
        'content': [
            "Objective: Interpret the model results and provide actionable insights.",
            "Activities:",
            "- Model Evaluation (Accuracy, Precision, Recall, F1 Score)",
            "- Identifying Key Predictors of High-Value Customers",
            "- Visualization of Results (e.g., feature importance, customer segmentation)",
            "- Recommendations for Targeting New Customers",
            "- Reporting and Documentation"
        ]
    }
}


# Function to add content to slides with checks for placeholder existence
def add_content_to_slide_checked(slide, title, content):
    # Check if title placeholder exists
    if slide.shapes.title:
        title_placeholder = slide.shapes.title
        title_placeholder.text = title
    
    # Check if content placeholder exists
    if len(slide.placeholders) > 1:
        content_placeholder = slide.placeholders[1]
        content_placeholder.text = '\n'.join(content)

# Add content to each slide
for i, (slide_title, slide_info) in enumerate(slides_content.items(), start=1):
    slide = presentation.slides[i]
    add_content_to_slide_checked(slide, slide_info['title'], slide_info['content'])

# Save the updated presentation
output_ppt_path = 'KPMG_Data_Analytics_Approach.pptx'
presentation.save(output_ppt_path)

output_ppt_path


# In[ ]:


# This presentation outlines the strategy behind each of the three phases—Data Exploration, Model Development, and Interpretation—detailing the activities involved in each phase. Let me know if you need any further modifications or additional information.

